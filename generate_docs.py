#!/usr/bin/env python3
"""
Generates PowerPoint documentation for College Management System
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx is not installed.")
    print("Please install it using: pip install python-pptx")
    exit(1)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title
    tf = body_shape.text_frame

    for i, point in enumerate(content_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    # Add title
    title_shape = shapes.title
    title_shape.text = title

    # Left column
    left = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_tf = left.text_frame
    left_tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            left_tf.text = point
        else:
            p = left_tf.add_paragraph()
            p.text = point

    # Right column
    right = shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    right_tf = right.text_frame
    right_tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            right_tf.text = point
        else:
            p = right_tf.add_paragraph()
            p.text = point

    return slide

def create_presentation():
    """Create the main presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                   "College Management System",
                   "Technical Documentation\nSpring Boot REST API with PostgreSQL")

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "A RESTful API for managing college operations",
        "Built with Spring Boot 4.0.1 and Java 17",
        "PostgreSQL database with Flyway migrations",
        "Swagger UI for API documentation",
        "Manages Students, Courses, Enrollments, and Payments",
        "Production-ready CRUD operations with validation"
    ])

    # Slide 3: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "Backend Framework: Spring Boot 4.0.1",
        "Language: Java 17",
        "Database: PostgreSQL",
        "ORM: Spring Data JPA / Hibernate",
        "Build Tool: Maven",
        "API Documentation: SpringDoc OpenAPI 3 (Swagger UI)",
        "Validation: Jakarta Validation",
        "Code Generation: Lombok",
        "Port: 8082"
    ])

    # Slide 4: System Architecture
    add_content_slide(prs, "System Architecture - Layered Design", [
        "Controller Layer: REST endpoints handling HTTP requests",
        "Service Layer: Business logic and transaction management",
        "Repository Layer: Data access using Spring Data JPA",
        "Entity Layer: JPA entities mapped to database tables",
        "Configuration Layer: OpenAPI/Swagger setup",
        "",
        "Pattern: MVC (Model-View-Controller) with REST",
        "Database Schema: 'college' in PostgreSQL 'collegedb'"
    ])

    # Slide 5: Database Schema
    add_content_slide(prs, "Database Schema", [
        "Schema: college",
        "",
        "Tables:",
        "  • student (id, name, address, phone_number)",
        "  • course (id, title, code)",
        "  • enrollment (id, student_id, course_id, enrollment_date, enrollment_status)",
        "  • payment (id, enroll_id, student_phone_number, payment_date, payment_status, amount)",
        "",
        "Relationships:",
        "  • Student 1:N Enrollment",
        "  • Course 1:N Enrollment",
        "  • Enrollment 1:N Payment"
    ])

    # Slide 6: Entity - Student
    add_content_slide(prs, "Entity: Student", [
        "Table: college.student",
        "",
        "Fields:",
        "  • id (Long, Primary Key, Auto-generated)",
        "  • name (String, required)",
        "  • address (String)",
        "  • phoneNumber (String, required)",
        "",
        "Business Rules:",
        "  • Phone number stored as String (supports country codes)",
        "  • Used as unique identifier for payment lookup"
    ])

    # Slide 7: Entity - Course
    add_content_slide(prs, "Entity: Course", [
        "Table: college.course",
        "",
        "Fields:",
        "  • id (Long, Primary Key, Auto-generated)",
        "  • title (String, required)",
        "  • code (String, required, unique course identifier)",
        "",
        "Example:",
        "  • title: 'Mathematics'",
        "  • code: 'MATH-101'"
    ])

    # Slide 8: Entity - Enrollment
    add_content_slide(prs, "Entity: Enrollment", [
        "Table: college.enrollment",
        "",
        "Fields:",
        "  • id (Long, Primary Key)",
        "  • student (ManyToOne relationship)",
        "  • course (ManyToOne relationship)",
        "  • enrollmentDate (LocalDate, required)",
        "  • enrollmentStatus (Enum: PENDING, ACTIVE, COMPLETED, CANCELLED)",
        "",
        "Business Rules:",
        "  • New enrollments start as PENDING",
        "  • Status becomes ACTIVE after successful payment",
        "  • Status managed automatically by PaymentService"
    ])

    # Slide 9: Entity - Payment
    add_content_slide(prs, "Entity: Payment", [
        "Table: college.payment",
        "",
        "Fields:",
        "  • id (Long, Primary Key)",
        "  • enrollment (ManyToOne, required - foreign key)",
        "  • studentPhoneNumber (String, for search/compatibility)",
        "  • paymentDate (LocalDate)",
        "  • paymentStatus (Enum: PENDING, SUCCESS, FAILED)",
        "  • amount (Double, required)",
        "",
        "Business Logic:",
        "  • Payment SUCCESS triggers enrollment status → ACTIVE"
    ])

    # Slide 10: API - Student Endpoints
    add_content_slide(prs, "API: Student Endpoints", [
        "Base Path: /students",
        "",
        "GET /students - List all students",
        "GET /students/{id} - Get student by ID",
        "GET /students/by-phone?phoneNumber={phone} - Get by phone",
        "POST /students - Create new student (201 Created)",
        "PUT /students/{id} - Update student",
        "DELETE /students/{id} - Delete student (204 No Content)",
        "",
        "Request Body: { name, address, phoneNumber }"
    ])

    # Slide 11: API - Course Endpoints
    add_content_slide(prs, "API: Course Endpoints", [
        "Base Path: /courses",
        "",
        "GET /courses - List all courses",
        "GET /courses/{id} - Get course by ID",
        "POST /courses - Create new course (201 Created)",
        "PUT /courses/{id} - Update course",
        "DELETE /courses/{id} - Delete course (204 No Content)",
        "",
        "Request Body: { title, code }",
        "",
        "Example: { \"title\": \"Mathematics\", \"code\": \"MATH-101\" }"
    ])

    # Slide 12: API - Enrollment Endpoints
    add_content_slide(prs, "API: Enrollment Endpoints", [
        "Base Path: /enrollments",
        "",
        "GET /enrollments - List all enrollments",
        "GET /enrollments/{id} - Get enrollment by ID",
        "POST /enrollments - Create enrollment (201 Created)",
        "PUT /enrollments/{id} - Update enrollment",
        "DELETE /enrollments/{id} - Delete enrollment (204 No Content)",
        "",
        "Request Body: { studentId, courseId, enrollmentDate }",
        "  • enrollmentDate defaults to today if omitted",
        "  • Initial status: PENDING (auto-managed)"
    ])

    # Slide 13: API - Payment Endpoints
    add_content_slide(prs, "API: Payment Endpoints", [
        "Base Path: /payments",
        "",
        "GET /payments - List all payments",
        "GET /payments/{id} - Get payment by ID",
        "POST /payments - Create payment (201 Created)",
        "PUT /payments/{id} - Update payment",
        "DELETE /payments/{id} - Delete payment (204 No Content)",
        "",
        "Request Body: { enrollmentId, paymentDate, amount }",
        "Response: { id, enrollmentId, studentPhoneNumber, paymentDate,",
        "           paymentStatus, amount }"
    ])

    # Slide 14: Business Logic Highlights
    add_content_slide(prs, "Business Logic Highlights", [
        "Enrollment Status Management:",
        "  • New enrollment → PENDING status",
        "  • Cannot manually set status via API",
        "  • Payment with SUCCESS status → Enrollment becomes ACTIVE",
        "",
        "Payment Processing:",
        "  • Must link to existing enrollment",
        "  • Auto-populates student phone from enrollment",
        "  • Payment date defaults to today if omitted",
        "",
        "Data Integrity:",
        "  • Foreign key validation on create/update",
        "  • Throws IllegalArgumentException for invalid references"
    ])

    # Slide 15: Configuration & Setup
    add_content_slide(prs, "Configuration & Setup", [
        "Database Prerequisites:",
        "  1. Install PostgreSQL",
        "  2. Create database: collegedb",
        "  3. Credentials (application.yml): postgres/postgres",
        "",
        "Application Configuration:",
        "  • Port: 8082",
        "  • Database: jdbc:postgresql://localhost:5432/collegedb",
        "  • Flyway: Auto-creates schema & tables on startup",
        "",
        "Run Application:",
        "  mvn spring-boot:run"
    ])

    # Slide 16: API Documentation
    add_content_slide(prs, "API Documentation & Testing", [
        "Swagger UI (Interactive):",
        "  http://localhost:8082/swagger-ui",
        "",
        "OpenAPI Specification:",
        "  http://localhost:8082/v3/api-docs",
        "",
        "Features:",
        "  • Try out API endpoints directly in browser",
        "  • View request/response schemas",
        "  • Auto-generated from code annotations",
        "  • Organized by tags: Students, Courses, Enrollments, Payments"
    ])

    # Slide 17: Project Structure
    add_content_slide(prs, "Project Structure", [
        "com.college.cms",
        "├── CollegeManagementApplication.java (Main class)",
        "├── config/",
        "│   ├── OpenApiConfig.java",
        "│   └── SwaggerTagOrderConfig.java",
        "├── controller/ (REST endpoints)",
        "│   ├── StudentController.java",
        "│   ├── CourseController.java",
        "│   ├── EnrollmentController.java",
        "│   └── PaymentController.java",
        "├── service/ (Business logic)",
        "├── repo/ (Data access - JpaRepository)",
        "└── entity/ (JPA entities)"
    ])

    # Slide 18: Key Dependencies
    add_content_slide(prs, "Key Maven Dependencies", [
        "spring-boot-starter-web - REST API support",
        "spring-boot-starter-data-jpa - Database ORM",
        "spring-boot-starter-validation - Input validation",
        "postgresql - PostgreSQL JDBC driver",
        "springdoc-openapi-starter-webmvc-ui (2.8.3) - Swagger",
        "lombok - Code generation (@Data, @Builder, etc.)",
        "",
        "Build:",
        "  • Maven",
        "  • Java 17",
        "  • Spring Boot 4.0.1"
    ])

    # Slide 19: Security & Best Practices
    add_content_slide(prs, "Design Patterns & Best Practices", [
        "Separation of Concerns:",
        "  • Controllers: Handle HTTP, validation",
        "  • Services: Business logic, transactions",
        "  • Repositories: Data access only",
        "",
        "Data Validation:",
        "  • @Valid annotations on request bodies",
        "  • @NotBlank, @NotNull constraints",
        "  • Custom validation in service layer",
        "",
        "API Design:",
        "  • RESTful principles",
        "  • Proper HTTP status codes (201, 204, etc.)",
        "  • Descriptive error messages"
    ])

    # Slide 20: Future Enhancements
    add_content_slide(prs, "Potential Future Enhancements", [
        "Security:",
        "  • Spring Security for authentication/authorization",
        "  • JWT token-based API access",
        "",
        "Features:",
        "  • Pagination for list endpoints",
        "  • Search and filtering capabilities",
        "  • Email notifications for enrollment/payment",
        "  • Grade management module",
        "  • Attendance tracking",
        "",
        "Technical:",
        "  • Caching (Redis)",
        "  • API rate limiting",
        "  • Comprehensive unit & integration tests"
    ])

    # Slide 21: Contact & Resources
    add_content_slide(prs, "Resources & Documentation", [
        "Project Files:",
        "  • README.md - Setup instructions",
        "  • pom.xml - Maven configuration",
        "",
        "API Documentation:",
        "  • Swagger UI: http://localhost:8082/swagger-ui",
        "  • OpenAPI Spec: http://localhost:8082/v3/api-docs",
        "",
        "Technology References:",
        "  • Spring Boot: https://spring.io/projects/spring-boot",
        "  • SpringDoc OpenAPI: https://springdoc.org",
        "  • PostgreSQL: https://www.postgresql.org"
    ])

    return prs

def main():
    """Main function"""
    print("Generating College Management System Documentation...")

    try:
        prs = create_presentation()
        output_file = "College_Management_System_Documentation.pptx"
        prs.save(output_file)
        print(f"✓ Successfully created: {output_file}")
        print(f"  Total slides: {len(prs.slides)}")

    except Exception as e:
        print(f"✗ Error creating presentation: {e}")
        return 1

    return 0

if __name__ == "__main__":
    exit(main())
